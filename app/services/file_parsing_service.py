import logging
import io
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Any

import docx
import pandas as pd
from PyPDF2 import PdfReader
from pptx import Presentation
from fastapi import UploadFile

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".xls", ".xlsx", ".csv"}


class ExtractionError(Exception):
    pass


@dataclass
class FileMetadata:
    # Always present
    filename: str
    extension: str
    file_size_bytes: int

    # Content shape
    word_count: int = 0
    char_count: int = 0

    # Format-specific counts
    page_count: int | None = None      # PDF
    slide_count: int | None = None     # PPTX
    sheet_names: list[str] = field(default_factory=list)   # XLS/XLSX
    paragraph_count: int | None = None # DOCX

    # Document properties (best-effort — not all formats expose these)
    title: str | None = None
    author: str | None = None
    created_at: datetime | None = None
    modified_at: datetime | None = None
    subject: str | None = None
    keywords: list[str] = field(default_factory=list)

    # Structural flags — useful for retrieval scoring
    has_tables: bool = False
    has_images: bool = False    # PPTX / DOCX only

    # Anything format-specific that doesn't fit above
    extra: dict[str, Any] = field(default_factory=dict)


@dataclass
class ExtractionResult:
    text: str
    metadata: FileMetadata




def extract_text_from_file(file: UploadFile) -> ExtractionResult:
    content = file.file.read()
    ext = Path(file.filename).suffix.lower()

    # print(f"Extracting text from '{file.filename}' (size={len(content)} bytes, ext='{ext}')")

    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(
            f"Unsupported file extension '{ext}'. "
            f"Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )

    try:
        result = _EXTRACTORS[ext](content, file.filename)
    except ExtractionError:
        raise
    except Exception as e:
        logger.exception("Failed to extract text from '%s'", file.filename)
        raise ExtractionError(f"Could not extract text from '{file.filename}': {e}")

    # Compute content-level metadata from the final text
    text = result.text.strip()
    result.metadata.char_count = len(text)
    result.metadata.word_count = len(text.split())
    result.metadata.file_size_bytes = len(content)

    return ExtractionResult(text=text, metadata=result.metadata)



def _extract_pdf(content: bytes, filename: str) -> ExtractionResult:
    reader = PdfReader(io.BytesIO(content))

    if reader.is_encrypted:
        raise ExtractionError("PDF is encrypted and cannot be read")


    raw = reader.metadata or {}

    def _pdf_date(val: str | None) -> datetime | None:
        """Parse PDF date string: D:YYYYMMDDHHmmSS"""
        if not val:
            return None
        try:
            cleaned = val.replace("D:", "")[:14]
            return datetime.strptime(cleaned, "%Y%m%d%H%M%S")
        except (ValueError, TypeError):
            return None
        

    meta = FileMetadata(
        filename=filename,
        extension=".pdf",
        file_size_bytes=0,
        page_count=len(reader.pages),
        keywords=[k.strip() for k in (raw.get("/Keywords") or "").split(",") if k.strip()],
        created_at=_pdf_date(raw.get("/CreationDate")),
        extra={k: v for k, v in raw.items() if k not in
               {"/Title", "/Author", "/Subject", "/Keywords", "/CreationDate", "/ModDate"}},
    )

    # print(meta)    

    pages = []
    for i, page in enumerate(reader.pages):
        try:
            pages.append(page.extract_text() or "")
        except Exception as e:
            logger.warning("Skipping PDF page %d: %s", i, e)

    return ExtractionResult(text="\n".join(pages), metadata=meta)


def _extract_docx(content: bytes, filename: str) -> ExtractionResult:
    doc = docx.Document(io.BytesIO(content))
    props = doc.core_properties

    meta = FileMetadata(
        filename=filename,
        extension=".docx",
        file_size_bytes=0,
        paragraph_count=len(doc.paragraphs),
        keywords=[k.strip() for k in (props.keywords or "").split(",") if k.strip()],
        created_at=props.created,
        has_tables=len(doc.tables) > 0,
        has_images=any(
            "image" in (r.part.content_type or "")
            for p in doc.paragraphs
            for r in p.runs
            if hasattr(r, "part") and r.part is not "None"
        ),
        extra={"revision": props.revision, "category": props.category or None},
    )

    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    return ExtractionResult(text=text, metadata=meta)


def _extract_pptx(content: bytes, filename: str) -> ExtractionResult:
    prs = Presentation(io.BytesIO(content))
    props = prs.core_properties

    has_images = any(
        shape.shape_type == 13  
        for slide in prs.slides
        for shape in slide.shapes
    )
    has_tables = any(
        shape.has_table
        for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, "has_table")
    )

    meta = FileMetadata(
        filename=filename,
        extension=".pptx",
        file_size_bytes=0,
        slide_count=len(prs.slides),
        title=props.title or 'None',
        author=props.author or 'None',
        subject=props.subject or 'None',
        keywords=[k.strip() for k in (props.keywords or "").split(",") if k.strip()],
        created_at=props.created,
        modified_at=props.modified,
        has_images=has_images,
        has_tables=has_tables,
    )

    chunks = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = [
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ]
        if texts:
            chunks.append(f"[Slide {i}]\n" + "\n".join(texts))

    return ExtractionResult(text="\n\n".join(chunks), metadata=meta)


def _extract_excel(content: bytes, filename: str) -> ExtractionResult:
    if filename.lower().endswith(".csv"):
        df = pd.read_csv(io.BytesIO(content))
        meta = FileMetadata(
            filename=filename,
            extension=".csv",
            file_size_bytes=0,
            has_tables=True,
            extra={"row_count": len(df), "column_count": len(df.columns)},
        )
        lines = df.fillna("").astype(str).agg(" | ".join, axis=1).tolist()
        return ExtractionResult(text="\n".join(lines), metadata=meta)

    xls = pd.ExcelFile(io.BytesIO(content))
    ext = Path(filename).suffix.lower()

    meta = FileMetadata(
        filename=filename,
        extension=ext,
        file_size_bytes=0,
        sheet_names=xls.sheet_names,
        has_tables=True, 
        extra={"sheet_count": len(xls.sheet_names)},
    )

    chunks = []
    for sheet in xls.sheet_names:
        df = xls.parse(sheet)
        if df.empty:
            continue
        rows = df.fillna("").astype(str)
        rows = rows[rows.apply(lambda r: r.str.strip().any(), axis=1)]
        lines = rows.agg(" | ".join, axis=1).tolist()
        if lines:
            chunks.append(f"[Sheet: {sheet}]\n" + "\n".join(lines))

    return ExtractionResult(text="\n\n".join(chunks), metadata=meta)




_EXTRACTORS = {
    ".pdf":  _extract_pdf,
    ".docx": _extract_docx,
    ".pptx": _extract_pptx,
    ".xls":  _extract_excel,
    ".xlsx": _extract_excel,
}